"""
Enterprise Agent Platform v3 - PowerPoint Presentation Generator
Author: Jan Střechovský (Honza)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ─── Color Palette ────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1B, 0x36, 0x5D)   # #1B365D  — primary
TEAL        = RGBColor(0x00, 0x97, 0xA7)   # #0097A7  — accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF4, 0xF6, 0xF8)
MID_GRAY    = RGBColor(0xB0, 0xBE, 0xC5)
DARK_GRAY   = RGBColor(0x37, 0x47, 0x4F)
LIGHT_TEAL  = RGBColor(0xE0, 0xF7, 0xFA)
AMBER       = RGBColor(0xFF, 0xB3, 0x00)   # warning accent
GREEN       = RGBColor(0x2E, 0x7D, 0x32)   # success

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── Helpers ──────────────────────────────────────────────────────────────────

def hex_rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def set_cell_bg(cell, rgb):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    srgbClr = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgbClr.set('val', '%02X%02X%02X' % (rgb[0], rgb[1], rgb[2]))


def tf_paragraph(tf, text, bold=False, size=Pt(14), color=WHITE, align=PP_ALIGN.LEFT, italic=False, space_before=None):
    p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = size
    run.font.color.rgb = color
    run.font.italic = italic
    if space_before:
        p.space_before = space_before
    return p


def clear_tf(tf):
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    tf.paragraphs[0].clear()


def set_tf_text(tf, text, bold=False, size=Pt(14), color=WHITE, align=PP_ALIGN.LEFT):
    clear_tf(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = size
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text="", bold=False,
                size=Pt(14), color=WHITE, align=PP_ALIGN.LEFT,
                word_wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    clear_tf(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = size
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_slide_number(slide, slide_num, total_slides=33):
    num_box = slide.shapes.add_textbox(
        Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.3)
    )
    tf = num_box.text_frame
    clear_tf(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{slide_num}"
    run.font.size = Pt(9)
    run.font.color.rgb = MID_GRAY


def full_bg(slide, color):
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)
    bg.z_order = 0


def dark_header_band(slide, height=Inches(1.5)):
    band = add_rect(slide, 0, 0, SLIDE_W, height, DARK_BLUE)
    return band


def bottom_bar(slide, color=TEAL, height=Inches(0.08)):
    bar = add_rect(slide, 0, SLIDE_H - height, SLIDE_W, height, color)
    return bar


def left_accent_bar(slide, color=TEAL, width=Inches(0.12)):
    bar = add_rect(slide, 0, 0, width, SLIDE_H, color)
    return bar


# ─── Slide Templates ──────────────────────────────────────────────────────────

def make_title_slide(prs, title, subtitle, author):
    """Full bleed dark title slide."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Full dark blue background
    full_bg(slide, DARK_BLUE)

    # Teal accent left bar
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, TEAL)

    # Teal accent bottom bar
    add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), TEAL)

    # Decorative teal rectangle (right side)
    add_rect(slide, Inches(10.5), Inches(1.5), Inches(2.5), Inches(4), RGBColor(0x00, 0x7A, 0x87))

    # Large white diagonal decorative shape suggestion via rectangle
    add_rect(slide, Inches(9.8), Inches(0), Inches(0.08), SLIDE_H, TEAL)

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(9.0), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(44)
    run.font.color.rgb = WHITE

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.6), Inches(9.0), Inches(1.0))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    clear_tf(tf2)
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(20)
    run2.font.color.rgb = TEAL
    run2.font.bold = False

    # Author
    auth_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(6.0), Inches(0.5))
    tf3 = auth_box.text_frame
    clear_tf(tf3)
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = f"Presented by: {author}"
    run3.font.size = Pt(13)
    run3.font.color.rgb = MID_GRAY

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.1), Inches(6.0), Inches(0.4))
    tf4 = date_box.text_frame
    clear_tf(tf4)
    p4 = tf4.paragraphs[0]
    run4 = p4.add_run()
    run4.text = "February 2026  |  Confidential & Internal Use Only"
    run4.font.size = Pt(11)
    run4.font.color.rgb = MID_GRAY
    run4.font.italic = True

    return slide


def make_section_divider(prs, section_num, section_title, subtitle=""):
    """Section break slide — teal left half, dark blue right."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Left teal block
    add_rect(slide, 0, 0, Inches(5.5), SLIDE_H, TEAL)

    # Right dark blue block
    add_rect(slide, Inches(5.5), 0, Inches(7.83), SLIDE_H, DARK_BLUE)

    # Section number large
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.0))
    tf = num_box.text_frame
    clear_tf(tf)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{section_num:02d}"
    run.font.bold = True
    run.font.size = Pt(96)
    run.font.color.rgb = WHITE

    # Section label
    lbl_box = slide.shapes.add_textbox(Inches(5.8), Inches(2.5), Inches(7.0), Inches(0.6))
    tf2 = lbl_box.text_frame
    clear_tf(tf2)
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = f"Section {section_num}"
    run2.font.size = Pt(14)
    run2.font.color.rgb = TEAL
    run2.font.bold = True

    # Section title
    title_box = slide.shapes.add_textbox(Inches(5.8), Inches(3.1), Inches(7.0), Inches(1.8))
    tf3 = title_box.text_frame
    tf3.word_wrap = True
    clear_tf(tf3)
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = section_title
    run3.font.bold = True
    run3.font.size = Pt(30)
    run3.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(5.8), Inches(5.0), Inches(7.0), Inches(0.8))
        tf4 = sub_box.text_frame
        tf4.word_wrap = True
        clear_tf(tf4)
        p4 = tf4.paragraphs[0]
        run4 = p4.add_run()
        run4.text = subtitle
        run4.font.size = Pt(14)
        run4.font.color.rgb = MID_GRAY

    return slide


def make_content_slide(prs, title, bullets, slide_num,
                       two_col=False, col1_title="", col2_title="",
                       col1_bullets=None, col2_bullets=None,
                       icon_bullets=None):
    """Standard content slide with header band."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    full_bg(slide, WHITE)

    # Dark blue header band
    dark_header_band(slide, Inches(1.35))

    # Left teal accent bar
    add_rect(slide, 0, 0, Inches(0.08), SLIDE_H, TEAL)

    # Bottom teal bar
    bottom_bar(slide)

    # Slide title in header
    title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.2), Inches(12.5), Inches(0.95))
    tf = title_box.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = WHITE

    add_slide_number(slide, slide_num)

    if two_col and col1_bullets and col2_bullets:
        # Column 1
        if col1_title:
            c1t = slide.shapes.add_textbox(Inches(0.35), Inches(1.55), Inches(6.0), Inches(0.5))
            set_tf_text(c1t.text_frame, col1_title, bold=True, size=Pt(14), color=DARK_BLUE)
        c1 = slide.shapes.add_textbox(Inches(0.35), Inches(2.1), Inches(6.0), Inches(4.8))
        c1.text_frame.word_wrap = True
        clear_tf(c1.text_frame)
        for i, b in enumerate(col1_bullets):
            if i == 0:
                p = c1.text_frame.paragraphs[0]
            else:
                p = c1.text_frame.add_paragraph()
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = b
            run.font.size = Pt(13)
            run.font.color.rgb = DARK_GRAY

        # Column 2
        if col2_title:
            c2t = slide.shapes.add_textbox(Inches(6.8), Inches(1.55), Inches(6.0), Inches(0.5))
            set_tf_text(c2t.text_frame, col2_title, bold=True, size=Pt(14), color=DARK_BLUE)
        c2 = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(6.0), Inches(4.8))
        c2.text_frame.word_wrap = True
        clear_tf(c2.text_frame)
        for i, b in enumerate(col2_bullets):
            if i == 0:
                p = c2.text_frame.paragraphs[0]
            else:
                p = c2.text_frame.add_paragraph()
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = b
            run.font.size = Pt(13)
            run.font.color.rgb = DARK_GRAY

        # Divider line
        add_rect(slide, Inches(6.55), Inches(1.5), Inches(0.02), Inches(5.7), LIGHT_GRAY)

    elif bullets:
        bul_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(12.5), Inches(5.7))
        bul_box.text_frame.word_wrap = True
        clear_tf(bul_box.text_frame)
        for i, b in enumerate(bullets):
            if i == 0:
                p = bul_box.text_frame.paragraphs[0]
            else:
                p = bul_box.text_frame.add_paragraph()
            p.space_before = Pt(8)
            run = p.add_run()
            run.text = b
            if b.startswith("   ") or b.startswith("\t"):
                run.font.size = Pt(12)
                run.font.color.rgb = DARK_GRAY
            elif b.isupper() or (len(b) > 2 and b[0].isupper() and b[1] == " " and not b.startswith("  ")):
                run.font.size = Pt(13.5)
                run.font.color.rgb = DARK_GRAY
            else:
                run.font.size = Pt(13.5)
                run.font.color.rgb = DARK_GRAY

    return slide


def make_big_stat_slide(prs, title, stats, slide_num):
    """Slide with large metric boxes."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    full_bg(slide, WHITE)
    dark_header_band(slide, Inches(1.35))
    add_rect(slide, 0, 0, Inches(0.08), SLIDE_H, TEAL)
    bottom_bar(slide)

    title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.2), Inches(12.5), Inches(0.95))
    tf = title_box.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = WHITE
    add_slide_number(slide, slide_num)

    n = len(stats)
    box_w = Inches(12.5) / n
    gap = Inches(0.15)
    for i, (number, label, sub) in enumerate(stats):
        x = Inches(0.35) + i * (box_w + gap / n)
        box = add_rect(slide, x, Inches(1.7), box_w - gap, Inches(4.8),
                       DARK_BLUE if i % 2 == 0 else RGBColor(0x0E, 0x4C, 0x78))
        num_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.2), box_w - gap - Inches(0.2), Inches(1.8))
        set_tf_text(num_box.text_frame, number, bold=True, size=Pt(40), color=TEAL, align=PP_ALIGN.CENTER)
        lbl_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(4.1), box_w - gap - Inches(0.2), Inches(1.0))
        set_tf_text(lbl_box.text_frame, label, bold=True, size=Pt(14), color=WHITE, align=PP_ALIGN.CENTER)
        if sub:
            sub_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(5.1), box_w - gap - Inches(0.2), Inches(1.0))
            sub_box.text_frame.word_wrap = True
            set_tf_text(sub_box.text_frame, sub, size=Pt(11), color=MID_GRAY, align=PP_ALIGN.CENTER)

    return slide


def make_comparison_table(prs, title, headers, rows, slide_num):
    """Comparison table slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    full_bg(slide, WHITE)
    dark_header_band(slide, Inches(1.35))
    add_rect(slide, 0, 0, Inches(0.08), SLIDE_H, TEAL)
    bottom_bar(slide)

    title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.2), Inches(12.5), Inches(0.95))
    tf = title_box.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = WHITE
    add_slide_number(slide, slide_num)

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    table_left = Inches(0.35)
    table_top = Inches(1.55)
    table_width = Inches(12.6)
    table_height = Inches(5.7)

    table = slide.shapes.add_table(n_rows, n_cols, table_left, table_top,
                                   table_width, table_height).table

    col_widths = [Inches(3.5)] + [Inches(3.0)] * (n_cols - 1)
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        set_cell_bg(cell, (0x1B, 0x36, 0x5D))
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.text = header
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = WHITE

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            if ri % 2 == 0:
                set_cell_bg(cell, (0xF4, 0xF6, 0xF8))
            else:
                set_cell_bg(cell, (0xFF, 0xFF, 0xFF))
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            # Clear existing runs
            for run in p.runs:
                run.text = ""
            run = p.add_run() if not p.runs else p.runs[0]
            run.text = str(val)
            run.font.size = Pt(12)
            if ci == 0:
                run.font.bold = True
                run.font.color.rgb = DARK_BLUE
            elif val in ("Yes", "YES", "Full", "Native", "On-Prem"):
                run.font.color.rgb = GREEN
                run.font.bold = True
            elif val in ("No", "NO", "None", "Limited", "N/A", "Cloud Only"):
                run.font.color.rgb = RGBColor(0xC6, 0x28, 0x28)
            elif val in ("Partial", "Basic", "Add-on"):
                run.font.color.rgb = AMBER
            else:
                run.font.color.rgb = DARK_GRAY

    return slide


def make_architecture_slide(prs, title, slide_num):
    """Architecture overview slide with layered boxes."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    full_bg(slide, WHITE)
    dark_header_band(slide, Inches(1.35))
    add_rect(slide, 0, 0, Inches(0.08), SLIDE_H, TEAL)
    bottom_bar(slide)

    title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.2), Inches(12.5), Inches(0.95))
    tf = title_box.text_frame
    tf.word_wrap = True
    clear_tf(tf)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = WHITE
    add_slide_number(slide, slide_num)

    layers = [
        ("PRESENTATION LAYER",    "React 19 SPA  |  REST / WebSocket  |  Responsive UI",                  TEAL,           WHITE),
        ("API GATEWAY",           "FastAPI + Uvicorn  |  OIDC/JWT Auth  |  Rate Limiting  |  CORS",        DARK_BLUE,      WHITE),
        ("AGENT RUNTIME",         "Reasoning Engine  |  Tool Gateway  |  Memory Store  |  Plugin Host",    RGBColor(0x0E,0x4C,0x78), WHITE),
        ("LLM ROUTER",            "3-Tier Model Routing  |  LIGHT / STANDARD / HEAVY  |  Cost Optimizer",  RGBColor(0x15,0x5E,0x75), WHITE),
        ("DATA LAYER",            "PostgreSQL 16 + pgvector  |  Redis 7  |  S3-Compatible Object Store",   DARK_GRAY,      WHITE),
        ("OBSERVABILITY",         "Prometheus + Grafana  |  Loki + Promtail  |  Audit Trail",              RGBColor(0x1A,0x23,0x7E), WHITE),
    ]

    box_h = Inches(0.78)
    gap = Inches(0.04)
    start_y = Inches(1.45)
    box_left = Inches(0.35)
    box_w = Inches(12.6)

    for i, (layer_name, detail, bg, fg) in enumerate(layers):
        y = start_y + i * (box_h + gap)
        add_rect(slide, box_left, y, box_w, box_h, bg)

        lbl_box = slide.shapes.add_textbox(box_left + Inches(0.15), y + Inches(0.1),
                                           Inches(2.8), box_h - Inches(0.15))
        lbl_tf = lbl_box.text_frame
        clear_tf(lbl_tf)
        lp = lbl_tf.paragraphs[0]
        lr = lp.add_run()
        lr.text = layer_name
        lr.font.bold = True
        lr.font.size = Pt(10)
        lr.font.color.rgb = TEAL if bg != TEAL else WHITE

        det_box = slide.shapes.add_textbox(box_left + Inches(3.0), y + Inches(0.12),
                                           Inches(9.5), box_h - Inches(0.2))
        det_tf = det_box.text_frame
        clear_tf(det_tf)
        dp = det_tf.paragraphs[0]
        dr = dp.add_run()
        dr.text = detail
        dr.font.size = Pt(12)
        dr.font.color.rgb = WHITE

    return slide


def make_closing_slide(prs, slide_num):
    """Q&A / Thank you slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    full_bg(slide, DARK_BLUE)
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, TEAL)
    add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), TEAL)

    # Large Q&A text
    qa_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.0), Inches(2.5))
    tf = qa_box.text_frame
    clear_tf(tf)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Questions & Discussion"
    run.font.bold = True
    run.font.size = Pt(42)
    run.font.color.rgb = WHITE

    ty_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11.0), Inches(1.0))
    tf2 = ty_box.text_frame
    clear_tf(tf2)
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Thank you for your attention."
    run2.font.size = Pt(20)
    run2.font.color.rgb = TEAL

    contact_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.0), Inches(1.0))
    tf3 = contact_box.text_frame
    clear_tf(tf3)
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "Jan Strechovský (Honza)  |  Enterprise Agent Platform v3  |  February 2026"
    run3.font.size = Pt(13)
    run3.font.color.rgb = MID_GRAY
    run3.font.italic = True

    add_slide_number(slide, slide_num)
    return slide


# ─── Main Build ───────────────────────────────────────────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_num = 0

    # ── SLIDE 1 – Title ──────────────────────────────────────────────────────
    slide_num += 1
    make_title_slide(
        prs,
        title="Enterprise Agent Platform",
        subtitle="On-Premise AI Agent Infrastructure for Manufacturing Excellence",
        author="Jan Strechovský (Honza)",
    )

    # ══ SECTION 1 ════════════════════════════════════════════════════════════
    slide_num += 1
    make_section_divider(prs, 1, "Executive Summary",
                         "The business case for on-premise AI")

    # ── SLIDE 3 – The Challenge ───────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "The Challenge: Why AI — Why Now", [
        "  Manufacturing is undergoing the fastest transformation in a generation",
        "",
        "  Competitive Pressure",
        "    Competitors already use AI to cut cycle times and reduce defect rates",
        "    Manual data analysis cannot keep pace with production velocity",
        "",
        "  Operational Complexity",
        "    SAP ERP + MES + quality systems generate terabytes of unstructured data",
        "    Engineers spend 30-40 % of their time searching for information",
        "    Repetitive tasks consume expert capacity that should go to innovation",
        "",
        "  Compliance Burden",
        "    GDPR, ISO 27001, SOC 2, EAR / ITAR export controls require tight data governance",
        "    Any AI tool that sends data to external servers is an immediate compliance violation",
        "",
        "  The Opportunity",
        "    A purpose-built on-premise AI platform unlocks productivity while keeping full control",
    ], slide_num)

    # ── SLIDE 4 – The Problem with ChatGPT Enterprise ─────────────────────────
    slide_num += 1
    make_content_slide(prs, "The Problem with Cloud AI (ChatGPT Enterprise / Azure OpenAI)", [
        "  Data Sovereignty — Your data leaves the building",
        "    Every query, document, and business context sent to US-based servers",
        "    Zero guarantee that training opt-out is honoured end-to-end",
        "    GDPR Article 44 cross-border transfers — requires SCCs and DPAs",
        "",
        "  No Compliance Control",
        "    Cannot enforce ITAR / EAR export controls on cloud endpoints",
        "    Audit trail is whatever the vendor decides to provide",
        "    Cannot certify SOC 2 Type II or ISO 27001 against vendor infrastructure",
        "",
        "  No Deep Integration",
        "    ChatGPT has no access to your SAP or MES — copy-paste is the 'integration'",
        "    No real-time production data; stale context; hallucinations on internal processes",
        "",
        "  No Customization",
        "    Cannot fine-tune on your proprietary process data, domain vocabulary, or workflows",
        "    One-size-fits-all model — not trained for your quality standards or part numbers",
        "",
        "  Unpredictable Cost at Scale",
        "    Per-token pricing becomes prohibitive once usage scales across departments",
    ], slide_num)

    # ── SLIDE 5 – Our Solution ────────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Our Solution: Enterprise Agent Platform v3", [
        "  A fully on-premise, production-grade AI agent platform built specifically for",
        "  our manufacturing and enterprise environment.",
        "",
        "  Data never leaves our datacenter — Zero Egress Architecture",
        "",
        "  Agents that reason, verify, and act — with human approval gates",
        "",
        "  Direct integration with SAP and MES — no intermediaries",
        "",
        "  Compliance-by-design: SOC 2, GDPR, ISO 27001, EAR/ITAR",
        "",
        "  Multi-tenant: each department gets isolated, auditable access",
        "",
        "  207 source files, ~57,000 lines of production code",
        "  81 test files, ~27,000 lines of tests — enterprise-grade quality",
        "",
        "  Not a proof-of-concept.  This is production-ready infrastructure.",
        "  We need a server to run it.",
    ], slide_num)

    # ── SLIDE 6 – Key Benefits ────────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs,
        "Key Benefits Overview", [],
        slide_num,
        two_col=True,
        col1_title="For Management",
        col2_title="For IT / Engineering",
        col1_bullets=[
            "  Full data sovereignty — nothing leaves premises",
            "",
            "  Regulatory compliance out of the box",
            "  (GDPR, SOC 2, ISO 27001)",
            "",
            "  Cost predictability — fixed hardware vs",
            "  ever-growing per-token cloud bills",
            "",
            "  Fine-tuned on our data — higher accuracy",
            "  than generic cloud AI on our domain",
            "",
            "  SAP & MES integration — AI works with",
            "  real production data in real time",
            "",
            "  Human-in-the-Loop approval gates —",
            "  no automated writes without sign-off",
        ],
        col2_bullets=[
            "  FastAPI + PostgreSQL 16 + Redis 7 stack",
            "  — familiar, battle-tested components",
            "",
            "  Kubernetes-native: Helm charts, HPA, PDB,",
            "  NetworkPolicy — fits existing infra",
            "",
            "  Grafana + Prometheus + Loki — full",
            "  observability out of the box",
            "",
            "  Plugin architecture — extend without",
            "  modifying core platform",
            "",
            "  14 Alembic migrations — clean, versioned",
            "  database schema evolution",
            "",
            "  Load-tested with Locust + k6",
        ],
    )

    # ══ SECTION 2 ════════════════════════════════════════════════════════════
    slide_num += 1
    make_section_divider(prs, 2, "Why On-Premise?",
                         "The server purchase argument")

    # ── SLIDE 8 – Data Sovereignty ────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Data Sovereignty: Zero Data Egress Architecture", [
        "  Core Principle: Every byte stays inside our datacenter",
        "",
        "  What this means in practice:",
        "    LLM inference runs locally — queries never sent to OpenAI, Anthropic, or Azure",
        "    Documents, embeddings, and vector indices stored in our PostgreSQL instance",
        "    Redis cache is internal — no external calls for session or context data",
        "    All API traffic terminates at our Kubernetes ingress controller",
        "",
        "  Technical enforcement:",
        "    Kubernetes NetworkPolicy — egress rules block outbound to the internet",
        "    No API keys for external LLM services in the production configuration",
        "    PII redaction runs before any data touches the LLM — BLOCK/REDACT/WARN hierarchy",
        "    Data classification enforcement (PUBLIC / INTERNAL / CONFIDENTIAL / RESTRICTED)",
        "",
        "  Compliance impact:",
        "    GDPR Article 44 compliance is trivially satisfied — no cross-border transfer",
        "    ITAR/EAR controlled technical data never reaches foreign nationals via cloud",
        "    Audit trail of every data access event, exportable as JSON or CSV",
    ], slide_num)

    # ── SLIDE 9 – Cost Comparison ─────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Cost Analysis: On-Premise vs Cloud AI at Scale", [
        "  Cloud AI (ChatGPT Enterprise / Azure OpenAI) — usage-based pricing",
        "    ChatGPT Enterprise: ~$30 per user / month minimum",
        "    GPT-4o API: ~$5 per 1M input tokens, ~$15 per 1M output tokens",
        "    50 engineers x 500 queries/day x avg. 2K tokens = 50M tokens/day",
        "    Monthly cloud cost estimate: $75,000 – $150,000+ per month",
        "    Annual run rate: $900K – $1.8M — and growing with usage",
        "",
        "  On-Premise (our platform + one server)",
        "    Server hardware (GPU workstation / rack server): one-time ~$30,000 – $60,000",
        "    Electricity + cooling: ~$500 – $1,500 / month",
        "    Platform: already built — no additional licensing",
        "    Break-even vs cloud: 1 – 3 months at medium usage",
        "",
        "  Additional on-prem advantages:",
        "    3-tier model router uses LIGHT models (fast, cheap) for simple queries",
        "    Only escalates to HEAVY models when complexity requires it",
        "    Tenant budget caps enforced per department — no runaway spend",
        "    Fine-tuned models on our data = fewer tokens needed for same accuracy",
    ], slide_num)

    # ── SLIDE 10 – Compliance ─────────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Compliance: What Cloud AI Cannot Provide", [
        "  SOC 2 Type II",
        "    Requires auditability of all data access — we provide full audit trail",
        "    Cloud vendors provide their own SOC 2, not yours",
        "    Our platform: every agent action logged with user, timestamp, tenant, result",
        "",
        "  GDPR",
        "    Right to erasure: we can delete any user's data from our systems",
        "    Data minimisation: PII redaction before LLM processing",
        "    DPA with a US cloud vendor is not a substitute for on-premise control",
        "",
        "  ISO 27001",
        "    Information security management requires knowing where data lives",
        "    Our platform: data at rest encrypted, data in transit TLS 1.3",
        "    Network segmentation via Kubernetes NetworkPolicy",
        "",
        "  EAR / ITAR (Export Administration Regulations / International Traffic in Arms)",
        "    Controlled technical data (drawings, formulas, process specs) CANNOT go to cloud",
        "    Our export control tracking (TEC-06-02) classifies and guards this data",
        "    Data classification policy TEC-02-04: PUBLIC/INTERNAL/CONFIDENTIAL/RESTRICTED",
    ], slide_num)

    # ── SLIDE 11 – Customization ──────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Customization: What Generic Cloud AI Cannot Do", [
        "  Fine-tuning on our proprietary data",
        "    LoRA-ready fine-tuning pipeline — adapt base models to our terminology",
        "    Cloud AI is trained on the internet — not on our part numbers, quality specs, or processes",
        "    Result: fewer hallucinations, higher accuracy on domain-specific queries",
        "",
        "  Custom Reasoning Engine",
        "    OBSERVE → THINK → VERIFY structured reasoning loop",
        "    Agents can be configured with domain-specific tools and guardrails",
        "    Tool gateway controls what each agent is allowed to call",
        "",
        "  Custom Connectors",
        "    SAP/MES connector built in — read/write production data directly",
        "    Plugin architecture allows adding new connectors without touching core",
        "    No ChatGPT plugin can authenticate against our internal SAP system",
        "",
        "  Custom Workflows",
        "    Human-in-the-Loop (HITL) write approval — MFA-protected",
        "    Approval workflows mapped to our business processes, not generic templates",
        "    Multi-tenant: each department has isolated context, budget, and permissions",
    ], slide_num)

    # ── SLIDE 12 – SAP/MES Integration ───────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Integration: Direct SAP / MES Access", [
        "  Why integration matters",
        "    AI without real data is just a chatbot — useful for drafting emails, useless for ops",
        "    Our platform agents have direct, authenticated access to production systems",
        "",
        "  SAP Integration",
        "    Read: inventory levels, work orders, BOM, production orders, vendor master data",
        "    Write: create / update records — gated by Human-in-the-Loop approval + MFA",
        "    No copy-paste, no CSV export — live queries against SAP APIs",
        "",
        "  MES Integration",
        "    Read: machine status, OEE, production counts, quality measurements, alarms",
        "    Write: parameter adjustments — require HITL approval with audit trail",
        "    Real-time context for agents making recommendations",
        "",
        "  What this enables",
        "    Agent: 'Line 3 OEE dropped 8 % — root cause: tooling wear on station 7'",
        "    Agent: 'Propose maintenance work order in SAP' — engineer approves with MFA",
        "    This workflow is impossible with any cloud AI product today",
    ], slide_num)

    # ══ SECTION 3 ════════════════════════════════════════════════════════════
    slide_num += 1
    make_section_divider(prs, 3, "Architecture & Capabilities",
                         "Technical deep dive for the IT team")

    # ── SLIDE 14 – System Architecture ───────────────────────────────────────
    slide_num += 1
    make_architecture_slide(prs, "System Architecture Overview", slide_num)

    # ── SLIDE 15 – Multi-Tenant Security ─────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Multi-Tenant Security Architecture", [
        "  Tenant Isolation",
        "    Every database query filtered by tenant_id — no cross-tenant data leakage possible",
        "    Separate Redis namespaces per tenant",
        "    Tenant-level budget caps enforced on LLM usage",
        "",
        "  Identity & Access Management",
        "    OIDC / JWT: integrates with corporate identity providers (Azure AD, Keycloak)",
        "    SAML SSO support for enterprise single sign-on",
        "    API key authentication for service-to-service calls",
        "    Role-Based Access Control (RBAC) with fine-grained permissions",
        "",
        "  MFA for Critical Operations",
        "    TOTP-based multi-factor authentication required for write operations",
        "    Human-in-the-Loop approvals gated by MFA verification",
        "    Prevents automated or accidental writes to production systems",
        "",
        "  Session & Token Management",
        "    Short-lived JWT tokens with refresh rotation",
        "    Revocable API keys with scope restrictions",
        "    All auth events written to immutable audit trail",
    ], slide_num)

    # ── SLIDE 16 – AI Agent Runtime ───────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "AI Agent Runtime: How Agents Think & Act", [
        "  Structured Reasoning Engine",
        "    OBSERVE: collect context from tools, memory, and user input",
        "    THINK: multi-step reasoning with chain-of-thought logging",
        "    VERIFY: validate conclusions before proposing actions",
        "    This loop is deterministic and auditable — not a black box",
        "",
        "  Tool Gateway",
        "    Agents call external tools (SAP, MES, search, calculators) via controlled gateway",
        "    Each tool call is authenticated, rate-limited, and logged",
        "    Sandbox prevents agents from accessing unauthorized resources",
        "",
        "  Memory Architecture",
        "    Short-term: conversation context in Redis",
        "    Long-term: vectorised document memory in PostgreSQL + pgvector",
        "    Episodic: previous task outcomes stored for learning",
        "",
        "  Human-in-the-Loop (HITL)",
        "    Agents PROPOSE actions — humans APPROVE with MFA",
        "    No automated writes to SAP or MES without explicit sign-off",
        "    Full action plan shown to approver before execution",
    ], slide_num)

    # ── SLIDE 17 – RAG Pipeline ───────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "RAG Pipeline: Retrieval-Augmented Generation", [
        "  Document Ingestion",
        "    Upload PDFs, Word docs, text files, wiki pages",
        "    Chunking with configurable overlap and strategy",
        "    Metadata extraction and data classification tagging",
        "",
        "  Embedding & Storage",
        "    Dense vector embeddings stored in PostgreSQL with pgvector extension",
        "    Sparse BM25 index for keyword matching",
        "    Both stored on-premise — no external embedding API calls",
        "",
        "  Hybrid Search + Reranking",
        "    Combines vector similarity search with BM25 lexical search",
        "    Cross-encoder reranking improves result quality",
        "    Results filtered by tenant, classification, and user permissions",
        "",
        "  Query → Answer Flow",
        "    User question → hybrid retrieval → reranking → LLM generation",
        "    Sources cited in every answer — no hallucinated references",
        "    PII redaction applied before context is passed to LLM",
    ], slide_num)

    # ── SLIDE 18 – 3-Tier Model Router ────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "3-Tier LLM Model Router: Quality + Cost Optimisation", [
        "  LIGHT Tier — Fast & Cheap",
        "    Simple Q&A, classification, short summarisation",
        "    Sub-second latency, minimal compute",
        "    Handles ~60-70 % of all queries",
        "",
        "  STANDARD Tier — Balanced",
        "    Multi-step reasoning, document analysis, code generation",
        "    Moderate latency, medium compute",
        "    Handles ~25-30 % of queries",
        "",
        "  HEAVY Tier — Maximum Quality",
        "    Complex analysis, fine-tuned domain models, multi-agent orchestration",
        "    Higher latency, full GPU utilisation",
        "    Handles ~5-10 % of queries — only when justified",
        "",
        "  Routing Logic",
        "    Automatic routing based on query complexity scoring",
        "    Manual override available for power users",
        "    Budget caps prevent HEAVY tier abuse",
        "    Cost per query tracked per tenant and reported in Grafana",
    ], slide_num)

    # ── SLIDE 19 – HITL Workflow ──────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Human-in-the-Loop: Safe Agentic Operations", [
        "  The Problem it Solves",
        "    Autonomous AI writing to SAP without oversight is unacceptable",
        "    HITL makes AI useful AND safe — the agent does the work, humans stay in control",
        "",
        "  Workflow: PROPOSE → REVIEW → APPROVE → EXECUTE",
        "    1. Agent analyses data and formulates an action plan",
        "    2. Plan presented to designated approver with full context and reasoning",
        "    3. Approver reviews, can modify, must authenticate with MFA (TOTP)",
        "    4. Execution proceeds only after cryptographic approval",
        "    5. Full execution log written to audit trail",
        "",
        "  What Requires HITL",
        "    Any write operation to SAP (work orders, inventory, master data)",
        "    Any write operation to MES (parameter changes, alarms)",
        "    Sensitive data access above INTERNAL classification",
        "    Bulk operations affecting more than N records",
        "",
        "  Audit Completeness",
        "    Who proposed, who approved, what was executed, what changed",
        "    Exportable for SOC 2 audit in JSON and CSV format",
    ], slide_num)

    # ── SLIDE 20 – SAP/MES Connector ─────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "SAP / MES Connector Architecture", [
        "  Connector Framework",
        "    Typed connector interface — every integration implements the same contract",
        "    Connectors registered in plugin registry — hot-reload without restart",
        "    Per-connector authentication credentials stored in encrypted secrets",
        "",
        "  SAP Connector",
        "    RFC / REST / OData adapters depending on SAP version",
        "    Read operations: master data, transactional data, reporting",
        "    Write operations: work orders, goods movements, notifications",
        "    All writes go through HITL approval gate",
        "",
        "  MES Connector",
        "    OPC-UA / REST interface to shop-floor systems",
        "    Real-time machine data for agent context enrichment",
        "    Write-back for parameter adjustments (HITL gated)",
        "",
        "  Future Connectors (plugin architecture makes this easy)",
        "    Quality Management System (QMS)",
        "    Document Management System (DMS)",
        "    Maintenance Management",
        "    Supplier portals",
    ], slide_num)

    # ── SLIDE 21 – Plugin System ──────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Plugin Architecture: Built to Extend", [
        "  Design Philosophy",
        "    Core platform is stable — extensions via plugins, never by modifying core",
        "    Typed plugin interfaces — IDE autocompletion, compile-time safety",
        "    Sandboxed execution — a buggy plugin cannot crash the platform",
        "",
        "  Plugin Types",
        "    Tool Plugins: new capabilities for agents (calculators, file parsers, APIs)",
        "    Connector Plugins: new system integrations (ERP, MES, databases)",
        "    Model Plugins: swap LLM backends, add new model providers",
        "    Auth Plugins: custom authentication providers",
        "",
        "  Plugin Lifecycle",
        "    Discover → Validate → Register → Enable (hot-reload)",
        "    Disable and remove without restarting platform",
        "    Plugin SDK documented in PLUGIN_SDK.md",
        "",
        "  Security Sandboxing",
        "    Plugins declare required permissions at registration",
        "    Runtime permission enforcement — plugins cannot exceed declared scope",
        "    Plugin actions appear in audit trail with plugin identity",
    ], slide_num)

    # ══ SECTION 4 ════════════════════════════════════════════════════════════
    slide_num += 1
    make_section_divider(prs, 4, "Security & Compliance",
                         "Enterprise-grade protection by design")

    # ── SLIDE 23 – Security Architecture ─────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Security Architecture: Defence in Depth", [
        "  Authentication Layer",
        "    OIDC / JWT — integrates with Azure AD, Keycloak, Okta",
        "    SAML SSO — enterprise federation support",
        "    API keys — scoped service accounts for automation",
        "    MFA (TOTP) — required for all write operations",
        "",
        "  Authorisation Layer",
        "    RBAC with fine-grained permissions per resource type",
        "    Tenant isolation enforced at database query level",
        "    Rate limiting per user, per tenant, per endpoint",
        "",
        "  Transport Security",
        "    TLS 1.3 enforced for all in-cluster and external traffic",
        "    Kubernetes NetworkPolicy — default deny, explicit allow",
        "    Zero egress to internet — enforced at network level",
        "",
        "  Storage Security",
        "    PostgreSQL data-at-rest encryption",
        "    Secrets managed via Kubernetes Secrets / Vault integration",
        "    Encryption keys rotatable without downtime",
    ], slide_num)

    # ── SLIDE 24 – Data Protection ────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Data Protection: PII, Classification & Export Control", [
        "  PII Redaction (BLOCK / REDACT / WARN)",
        "    BLOCK: sensitive PII never enters LLM context (SSN, credentials, payment data)",
        "    REDACT: identifiers replaced with pseudonymous tokens before LLM processing",
        "    WARN: marginal PII flagged in audit trail for review",
        "    Regex-based engine — configurable patterns for new PII types",
        "",
        "  Data Classification (policy TEC-02-04)",
        "    PUBLIC — freely shareable, no restrictions",
        "    INTERNAL — employees only, standard controls",
        "    CONFIDENTIAL — restricted access, enhanced audit",
        "    RESTRICTED — HITL required for any access, executive sign-off",
        "",
        "  Export Control (EAR / ITAR — policy TEC-06-02)",
        "    Documents tagged with jurisdiction (EAR99, ECCN, ITAR)",
        "    Export control status checked before any sharing or agent use",
        "    Controlled technical data blocked from cloud, logged on access",
        "",
        "  Right to Erasure (GDPR Article 17)",
        "    User data deletion workflow — cascades across all tables",
        "    Audit trail of deletion events retained for compliance",
    ], slide_num)

    # ── SLIDE 25 – Audit Trail ────────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Audit Trail: Complete Accountability", [
        "  Every Action Logged",
        "    User authentication events (login, logout, token refresh, MFA challenges)",
        "    Every LLM query — user, tenant, model tier, prompt hash, token count",
        "    Every tool call — agent, tool, arguments, result, latency",
        "    Every HITL decision — proposer, approver, action, outcome",
        "    Every data access above INTERNAL classification",
        "",
        "  Audit Log Structure",
        "    Immutable append-only log (no update/delete via application layer)",
        "    Structured JSON with event type, actor, resource, timestamp, correlation ID",
        "    PostgreSQL-backed with retention policies",
        "",
        "  Export & Reporting",
        "    JSON export for SIEM integration",
        "    CSV export for SOC 2 auditor evidence packages",
        "    Compliance dashboard in Grafana — real-time posture view",
        "",
        "  Compliance Mapping",
        "    SOC 2 CC6, CC7, CC8 controls satisfied by audit trail",
        "    GDPR Article 30 Records of Processing Activities fulfilled",
        "    ISO 27001 Annex A.12.4 — logging and monitoring",
    ], slide_num)

    # ── SLIDE 26 – Network Security ───────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Network Security: Zero Egress, Zero Trust", [
        "  Kubernetes NetworkPolicy",
        "    Default-deny ingress and egress for all pods",
        "    Explicit allow rules for necessary service-to-service traffic only",
        "    Database pods: only accept connections from API pods",
        "    LLM inference pods: only accept from agent runtime pods",
        "",
        "  Zero Egress Architecture",
        "    No outbound internet from any production pod",
        "    Model weights downloaded and cached during controlled deployment",
        "    External connector calls (SAP, MES) via dedicated egress proxy with allow-list",
        "",
        "  Encryption in Transit",
        "    All inter-pod communication: mTLS via service mesh (Linkerd/Istio compatible)",
        "    Ingress: TLS 1.3 with HSTS",
        "    Database connections: SSL enforced",
        "",
        "  Encryption at Rest",
        "    PostgreSQL tablespace encryption",
        "    Redis persistence files encrypted on disk",
        "    PVC storage class with encryption enabled (LUKS / cloud-native)",
    ], slide_num)

    # ══ SECTION 5 ════════════════════════════════════════════════════════════
    slide_num += 1
    make_section_divider(prs, 5, "Operations & Reliability",
                         "Production-ready from day one")

    # ── SLIDE 28 – Monitoring ─────────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Monitoring & Observability Stack", [
        "  Metrics — Prometheus + Grafana",
        "    LLM Performance Dashboard: latency p50/p95/p99, token throughput, model tier distribution",
        "    Agent Operations Dashboard: task completion rate, tool call volume, error rates",
        "    Tenant Budget Dashboard: token spend per tenant, budget utilisation, cost trends",
        "    Infrastructure Dashboard: CPU/GPU/memory/disk, pod health, connection pool utilisation",
        "",
        "  Logging — Loki + Promtail",
        "    Structured JSON logs from all services",
        "    Log aggregation with label-based filtering (tenant, agent, model)",
        "    Correlation IDs link logs across service boundaries",
        "    Retention and archival policies configurable per log type",
        "",
        "  Alerting",
        "    Prometheus AlertManager rules for SLA breaches",
        "    PagerDuty / email / Slack integrations available",
        "    Budget overage alerts per tenant",
        "    Security event alerts (auth failures, rate limit spikes)",
        "",
        "  Application Performance",
        "    Distributed tracing ready (OpenTelemetry instrumentation)",
        "    Database query performance monitoring via pg_stat_statements",
    ], slide_num)

    # ── SLIDE 29 – Deployment ─────────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Deployment: From Development to Production", [
        "  Development Environment",
        "    Docker Compose stack — single command startup",
        "    All services (API, DB, Redis, workers) in containers",
        "    Hot-reload for rapid iteration",
        "    Matches production service topology",
        "",
        "  Production: Kubernetes via Helm",
        "    Helm charts with environment-specific values files",
        "    Horizontal Pod Autoscaler (HPA) — scale out under load",
        "    Pod Disruption Budget (PDB) — rolling updates with zero downtime",
        "    NetworkPolicy — security enforced at deploy time",
        "    Resource requests and limits set for every container",
        "",
        "  Multi-Region / High Availability",
        "    Multi-region failover support built into architecture",
        "    PostgreSQL streaming replication + automated failover",
        "    Redis Sentinel for cache HA",
        "    Stateless API layer scales horizontally without constraints",
        "",
        "  GitOps Ready",
        "    Helm values committed to git — full infrastructure as code",
        "    ArgoCD / Flux compatible for automated deployments",
    ], slide_num)

    # ── SLIDE 30 – Backup & DR ────────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Backup & Disaster Recovery", [
        "  PostgreSQL Backup Automation",
        "    Scheduled pg_dump via cron — configurable frequency",
        "    Backup verification: restore tested automatically after each backup",
        "    Retention policy: daily, weekly, monthly backups with automated purge",
        "    Backups encrypted and stored on separate storage volume",
        "",
        "  Recovery Procedures (documented in RUNBOOK.md)",
        "    RTO (Recovery Time Objective): < 30 minutes for full restore",
        "    RPO (Recovery Point Objective): < 1 hour with hourly backup schedule",
        "    Tested restore procedure — not just theoretical",
        "",
        "  Application State Recovery",
        "    Redis: persistence enabled (AOF + RDB), recovers from last checkpoint",
        "    Stateless API and agent services: restart immediately without data loss",
        "    14 Alembic migrations: schema recovery is deterministic and versioned",
        "",
        "  Incident Response",
        "    Runbook documented for all known failure scenarios",
        "    Grafana alerts trigger incident workflow",
        "    Immutable audit trail preserved even during incident",
    ], slide_num)

    # ── SLIDE 31 – Scaling ────────────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Scaling: Tested, Not Estimated", [
        "  Load Testing Suite",
        "    Locust scenarios: simulate N concurrent users with realistic query mix",
        "    k6 scripts: API throughput and latency under sustained load",
        "    LLM inference stress tests: token throughput at capacity",
        "",
        "  Horizontal Scaling (Kubernetes HPA)",
        "    API layer: scale to N replicas based on CPU + request rate",
        "    Worker pool: scale based on task queue depth",
        "    No shared mutable state in API layer — scales linearly",
        "",
        "  Database Scaling",
        "    PgBouncer connection pooling — handles thousands of connections",
        "    Read replicas for analytics and reporting workloads",
        "    pgvector HNSW index — sub-millisecond vector search at scale",
        "",
        "  Caching Strategy",
        "    Redis 7: LLM response caching for identical queries",
        "    Embedding cache: avoid re-embedding unchanged documents",
        "    Session cache: JWT validation without database round-trip",
    ], slide_num)

    # ══ SECTION 6 ════════════════════════════════════════════════════════════
    slide_num += 1
    make_section_divider(prs, 6, "Competitive Advantage",
                         "What we have that nobody else can offer you")

    # ── SLIDE 33 – Feature Comparison Table ───────────────────────────────────
    slide_num += 1
    headers = ["Feature", "Our Platform", "ChatGPT Enterprise", "Azure OpenAI"]
    rows = [
        ["Data stays on-premise",         "Yes",           "No",        "No"],
        ["GDPR / ITAR compliant",          "Yes",           "Partial",   "Partial"],
        ["SOC 2 audit trail (yours)",      "Yes",           "No",        "No"],
        ["SAP / MES integration",          "Native",        "No",        "No"],
        ["Fine-tuning on your data",       "Yes",           "No",        "Add-on"],
        ["Human-in-the-Loop workflow",     "Yes",           "No",        "No"],
        ["PII redaction before LLM",       "Yes",           "No",        "No"],
        ["Multi-tenant isolation",         "Yes",           "Limited",   "Partial"],
        ["Custom reasoning engine",        "Yes",           "No",        "No"],
        ["Plugin extensibility",           "Yes",           "Limited",   "Limited"],
        ["On-prem deployment",             "Yes",           "No",        "No"],
        ["Cost model",                     "Fixed (HW)",    "Per user",  "Per token"],
        ["Grafana / Prometheus monitoring","Yes",           "No",        "No"],
        ["Full source code ownership",     "Yes",           "No",        "No"],
    ]
    make_comparison_table(prs, "Feature Comparison: Our Platform vs Cloud AI", headers, rows, slide_num)

    # ── SLIDE 34 – What We Can Do That They Can't ─────────────────────────────
    slide_num += 1
    make_content_slide(prs, "What We Can Do That Cloud AI Cannot", [
        "  1. Query live SAP data and propose work orders — approved with MFA",
        "     No cloud AI product has authenticated access to your internal SAP",
        "",
        "  2. Fine-tune on confidential process data without it leaving the building",
        "     ITAR/EAR controlled data cannot legally go to cloud AI providers",
        "",
        "  3. Enforce your exact compliance policy at every query",
        "     PII redaction, data classification, export control — built in, not bolted on",
        "",
        "  4. Provide a complete, YOUR audit trail for SOC 2 / ISO 27001 evidence",
        "     Cloud vendors give you their audit — not yours",
        "",
        "  5. Scale cost-predictably with a fixed hardware investment",
        "     No surprise monthly bills when a department doubles its AI usage",
        "",
        "  6. Customise the reasoning engine and add proprietary tools",
        "     Open plugin architecture — extend for any future business need",
        "",
        "  7. Operate when internet connectivity is degraded or unavailable",
        "     Fully self-contained — no dependency on external APIs for core function",
    ], slide_num)

    # ══ SECTION 7 ════════════════════════════════════════════════════════════
    slide_num += 1
    make_section_divider(prs, 7, "Investment & Next Steps",
                         "What we need and what comes next")

    # ── SLIDE 36 – Server Requirements ───────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Server Requirements: Recommended Hardware Specification", [
        "  Minimum Configuration (pilot / small team, <20 concurrent users)",
        "    CPU: 16-core modern server CPU (AMD EPYC 7003 or Intel Xeon Scalable)",
        "    RAM: 128 GB ECC DDR4",
        "    GPU: 1x NVIDIA RTX 4090 24GB or A4000 16GB",
        "    Storage: 2 TB NVMe SSD (OS + DB) + 4 TB HDD (backups / model weights)",
        "    Estimated cost: ~$15,000 – $25,000",
        "",
        "  Recommended Configuration (production, <100 concurrent users)",
        "    CPU: 32-core server CPU",
        "    RAM: 256 GB ECC DDR4",
        "    GPU: 2x NVIDIA A100 40GB or H100 80GB",
        "    Storage: 4 TB NVMe NVMe (RAID1) + 20 TB NAS backup",
        "    Network: 25 GbE dual port NIC",
        "    Estimated cost: ~$40,000 – $80,000",
        "",
        "  Software Stack (all open-source / already licensed)",
        "    OS: Ubuntu Server 22.04 LTS",
        "    Container runtime: Docker 27 / containerd",
        "    Orchestration: Kubernetes 1.30 (k3s or full RKE2)",
        "    LLM backend: Ollama / vLLM / llama.cpp (model-dependent)",
    ], slide_num)

    # ── SLIDE 37 – Implementation Timeline ───────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Implementation Timeline: What Is Done vs What Comes Next", [
        "  COMPLETED — Platform Core",
        "    Architecture, multi-tenancy, auth (OIDC/SAML/MFA), RBAC",
        "    Agent runtime: reasoning engine, tool gateway, memory",
        "    RAG pipeline: ingestion, hybrid search, reranking",
        "    3-tier LLM router, HITL workflow, audit trail",
        "    SAP/MES connector framework, plugin system",
        "    Security: PII redaction, data classification, export control",
        "    Observability: Grafana, Prometheus, Loki",
        "    Kubernetes Helm charts, HPA, PDB, NetworkPolicy",
        "    Load testing suite (Locust + k6)",
        "",
        "  NEXT — After Server Procurement",
        "    Week 1-2: Server setup, OS, Kubernetes cluster provisioning",
        "    Week 3-4: LLM model download, fine-tuning baseline",
        "    Week 5-6: SAP connector configuration and UAT",
        "    Week 7-8: Pilot rollout to first department (5-10 users)",
        "    Month 3: Broader rollout + user training",
        "    Month 4-6: Fine-tuning on departmental data, advanced use cases",
    ], slide_num)

    # ── SLIDE 38 – ROI Projection ─────────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "ROI Projection: Expected Return on Investment", [
        "  Time Savings",
        "    Engineers spend avg. 2-3 hours/day searching for information",
        "    Estimated reduction: 1 hour/day per power user with AI-assisted search",
        "    50 engineers x 1 hour x 220 working days = 11,000 engineer-hours/year",
        "    At avg. loaded cost of 80 EUR/hr = 880,000 EUR/year in productivity",
        "",
        "  Error Reduction",
        "    SAP data entry errors: estimated 5-10 % of orders require correction",
        "    HITL-gated AI assistance: target 50 % reduction in entry errors",
        "    Each error correction: avg. 4 hours of engineering + admin time",
        "",
        "  Compliance Cost Avoidance",
        "    GDPR fine for a data breach involving cloud AI: up to 4 % annual revenue",
        "    SOC 2 audit cost reduction via automated evidence collection",
        "    On-premise = no costly data breach notification exercises",
        "",
        "  Hardware payback period",
        "    40,000 EUR server vs 75,000+ EUR/month cloud alternative",
        "    Break-even: < 1 month at target usage levels",
        "    3-year TCO advantage vs cloud: > 2,000,000 EUR",
    ], slide_num)

    # ── SLIDE 39 – Next Steps / CTA ───────────────────────────────────────────
    slide_num += 1
    make_content_slide(prs, "Next Steps: Decision Required", [
        "  What We Are Asking For",
        "    Approval to purchase on-premise server hardware",
        "    Recommended spec: 32-core CPU, 256 GB RAM, 2x GPU, 4 TB NVMe",
        "    Budget range: 40,000 – 80,000 EUR (one-time capital expenditure)",
        "    Rack space in our datacenter (2U – 4U)",
        "",
        "  What Happens Next (our commitment)",
        "    Week 1: Server spec finalisation and procurement order",
        "    Week 2-4: Hardware delivery, OS install, Kubernetes cluster setup",
        "    Week 5-6: Platform deployment, LLM model installation",
        "    Week 7-8: SAP connector UAT with IT team",
        "    Month 2: Pilot with first business unit (5-10 users)",
        "    Month 3: Go/no-go decision for broader rollout",
        "",
        "  Risks of NOT Proceeding",
        "    Continued manual processes while competitors automate",
        "    Compliance exposure if departments turn to unsanctioned cloud AI tools",
        "    Platform effort (57K lines of code, 81 test files) remains undeployed",
        "",
        "  We have the platform.  We need the server.  The ROI is clear.",
    ], slide_num)

    # ── SLIDE 40 – Q&A ────────────────────────────────────────────────────────
    slide_num += 1
    make_closing_slide(prs, slide_num)

    return prs


# ─── Entry Point ──────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import os
    out_path = "/home/honza/enterprise-agent-platform-v3/docs/Enterprise_Agent_Platform_Presentation.pptx"
    print(f"Building presentation...")
    prs = build_presentation()
    prs.save(out_path)
    size_kb = os.path.getsize(out_path) / 1024
    print(f"Saved: {out_path}")
    print(f"File size: {size_kb:.1f} KB")
    print(f"Slides: {len(prs.slides)}")
